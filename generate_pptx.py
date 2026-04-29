from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    content = slide.placeholders[1]
    content.text = content_text
    
    # Format title
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = 'Helvetica'
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(17, 17, 17)
        
    # Format content
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.name = 'Helvetica'
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(80, 80, 80)
        
    return slide

def create_catalog():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "EVO CORPORATE"
    subtitle.text = "Elegância em Movimento.\nExclusividade para a sua marca."
    
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(17, 17, 17)
        paragraph.font.name = 'Helvetica'
        
    # Slide 2: A EVO
    add_slide(prs, "ELEVANDO O PADRÃO DO SEU BRINDE", 
              "A EVO nasceu para unir design sofisticado, funcionalidade extrema e alta durabilidade. Especialistas em mochilas, malas e acessórios em couro 100% legítimo, criamos peças que acompanham rotinas exigentes com elegância.\n\n"
              "• Couro nobre de exportação (floater)\n"
              "• Design autoral e atemporal\n"
              "• Certificação LWG – garantia de origem sustentável")

    # Slide 3: Co-branding
    add_slide(prs, "SUA MARCA COM ASSINATURA DE EXCELÊNCIA",
              "O modelo Co-branding\n\n"
              "Oferecemos a possibilidade de personalizar nossos produtos com a sua logomarca, criando uma experiência única de pertencimento e luxo.\n\n"
              "• Gravação em baixo relevo ou clichê térmico diretamente no couro.\n"
              "• Aplicação sofisticada e minimalista.\n"
              "• Exemplo: Sua Marca by EVO.")

    # Slide 4: Mochila Vicenza
    add_slide(prs, "COLEÇÃO EXECUTIVA: MOCHILA VICENZA",
              "A escolha perfeita para o executivo dinâmico. Une praticidade e estilo para o dia a dia.\n\n"
              "Diferenciais:\n"
              "• Compartimento acolchoado para notebook de até 15,6”\n"
              "• Bolso traseiro antifurto com zíper oculto\n"
              "• Elástico traseiro para encaixe em malas de rodinhas\n"
              "• Capacidade: 12 Litros | Dimensões: 41 x 29 x 10 cm\n\n"
              "Cores Disponíveis: Preto, Café, Pinhão")

    # Slide 5: Mochila Selvaggio
    add_slide(prs, "COLEÇÃO EXECUTIVA: MOCHILA SELVAGGIO CLASSIC",
              "Um clássico atemporal. O equilíbrio perfeito entre a robustez do couro legítimo e o requinte necessário para ambientes corporativos.\n\n"
              "Diferenciais:\n"
              "• Estrutura reforçada 100% em couro legítimo\n"
              "• Acabamento interno premium\n"
              "• Alça de mão anatômica em couro\n\n"
              "Cores Disponíveis: Preto, Café, Pinhão")

    # Slide 6: Malas
    add_slide(prs, "COLEÇÃO VIAGEM: DUFFLE BAG EVO",
              "A companheira ideal para viagens curtas. Um presente corporativo de altíssimo impacto visual.\n\n"
              "Diferenciais:\n"
              "• Couro premium floater resistente e encorpado\n"
              "• Amplo espaço interno\n"
              "• Alça transversal removível e ajustável\n\n"
              "Cores Disponíveis: Preto, Café")

    # Slide 7: Kits
    add_slide(prs, "ACESSÓRIOS & KITS CORPORATIVOS",
              "Pequenos detalhes que transformam o ambiente de trabalho.\n\n"
              "• Necessarie EVO: Sofisticação para viagens.\n"
              "• Organizador de Mesa Dock: Elegância para o escritório.\n"
              "• Organizador Tech Dock: Para cabos e carregadores.\n"
              "• Deskpad DOCK: Conforto e proteção para a mesa.\n"
              "• Porta Cartão Clip EVO: Minimalismo para carregar o essencial.\n\n"
              "Sugestão de Kit: Deskpad + Organizador de Mesa + Porta Cartão")

    # Slide 8: Contato
    add_slide(prs, "VAMOS CONSTRUIR ESSA EXPERIÊNCIA JUNTOS?",
              "Nosso time está pronto para realizar um orçamento sob medida para o seu projeto B2B.\n\n"
              "Contatos Corporativos EVO:\n"
              "• WhatsApp: +55 19 98122-6232\n"
              "• E-mail: contato@useevo.store\n"
              "• Site: www.useevo.store\n\n"
              "EVO - ELEGÂNCIA EM MOVIMENTO LTDA.")

    prs.save('Catalogo_EVO_Corporate.pptx')
    print("PPTX gerado com sucesso!")

if __name__ == '__main__':
    create_catalog()
